# backend/files/processing_helpers.py
import os
import io
import re
import time
import base64
import logging
from typing import Optional, List, Tuple

from django.conf import settings

logger = logging.getLogger(__name__)

# Optional local OCR import - se mantiene como fallback
try:
    import easyocr
except Exception:
    easyocr = None

# Optional extractors (will be used if available)
try:
    from docx import Document
except Exception:
    Document = None

try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
except Exception:
    pdfminer_extract_text = None

# pdf -> images (optional)
try:
    from pdf2image import convert_from_path
except Exception:
    convert_from_path = None

# PIL for image handling
try:
    from PIL import Image
except Exception:
    Image = None

# OpenAI SDK import (repo appears to use "from openai import OpenAI")
try:
    from openai import OpenAI
except Exception:
    OpenAI = None

# Env/settings lookup helper
def _get_setting(name: str, default=None):
    return getattr(settings, name, os.getenv(name, default))

OPENAI_API_KEY = _get_setting("OPENAI_API_KEY", None)
OPENAI_MODEL_TEXT = _get_setting("OPENAI_MODEL_TEXT", "gpt-4o-mini")
OPENAI_MODEL_VISION = _get_setting("OPENAI_MODEL_VISION", "gpt-4o")
USE_LOCAL_OCR = _get_setting("USE_LOCAL_OCR", "False") in (True, "True", "true", "1")
MAX_SEND_SIZE_MB = int(_get_setting("MAX_SEND_SIZE_MB", 5))  # reject uploads larger than this by MB
CHUNK_SIZE_TOKENS = int(_get_setting("CHUNK_SIZE_TOKENS", 4000))  # approximate chunk size to send safely
MAX_RETRIES = int(_get_setting("OPENAI_MAX_RETRIES", 2))
DEFAULT_TEMPERATURE = float(_get_setting("OPENAI_TEMP", 0.0))
MAX_TOKENS_RESPONSE = int(_get_setting("OPENAI_MAX_TOKENS", 1500))

_openai_client = None

def get_openai_client():
    global _openai_client
    if _openai_client is not None:
        return _openai_client
    if OpenAI is None:
        raise RuntimeError("OpenAI SDK no está instalado en el entorno.")
    api_key = OPENAI_API_KEY or os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY no está configurada en settings ni en env.")
    _openai_client = OpenAI(api_key=api_key)
    return _openai_client

def text_to_md(text: str) -> str:
    """Limpieza básica y normalización para convertir texto plano a Markdown simple."""
    if not text:
        return ''
    # Normalize line endings
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    # Remove excessive blank lines
    text = re.sub(r'\n{3,}', '\n\n', text)
    # Strip leading/trailing whitespace
    return text.strip()

def _estimate_tokens_from_chars(chars: int) -> int:
    # heurística simple: ~4 chars per token (varía). Usada para chunking.
    return max(1, chars // 4)

def _chunk_text_by_chars(text: str, target_tokens: int = CHUNK_SIZE_TOKENS) -> List[str]:
    # conservatively convert tokens -> chars, using 4 chars per token
    approx_chars = target_tokens * 4
    if len(text) <= approx_chars:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = start + approx_chars
        # try to break at newline/space to avoid mid-sentence cuts
        if end < len(text):
            nl = text.rfind('\n', start, end)
            sp = text.rfind(' ', start, end)
            cut = nl if nl > start else sp if sp > start else end
        else:
            cut = len(text)
        chunks.append(text[start:cut].strip())
        start = cut
    return [c for c in chunks if c.strip()]

def _image_file_to_data_url_bytes(b: bytes, ext: str) -> str:
    ext = ext.lower().lstrip('.')
    mime = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{ext}"
    b64 = base64.b64encode(b).decode("utf-8")
    return f"data:{mime};base64,{b64}"

def _image_path_to_data_url(path: str) -> str:
    with open(path, "rb") as f:
        b = f.read()
    ext = os.path.splitext(path)[1].lower().lstrip('.')
    return _image_file_to_data_url_bytes(b, ext)

def _build_image_messages(prompt_text: str, data_url: str):
    """
    Construye mensajes para pasar imagen + prompt al endpoint de chat completions que soporte vision.
    Estructura pensada para robustez frente a variaciones de SDK.
    """
    user_content = [
        {"type": "text", "text": prompt_text},
        {"type": "image_url", "image_url": {"url": data_url}},
    ]
    return [{"role": "user", "content": user_content}]

def _call_openai_chat_completions(model: str, messages, max_tokens=MAX_TOKENS_RESPONSE, temperature=DEFAULT_TEMPERATURE, max_retries=MAX_RETRIES):
    client = get_openai_client()
    attempt = 0
    while True:
        try:
            resp = client.chat.completions.create(
                model=model,
                messages=messages,
                max_tokens=max_tokens,
                temperature=temperature,
            )
            # intentar extraer texto de respuesta robustamente
            if resp and getattr(resp, "choices", None):
                choice = resp.choices[0]
                out = None
                # soporte estructuras dict-like o objeto SDK
                try:
                    if isinstance(choice, dict):
                        msg = choice.get("message") or {}
                        out = msg.get("content") or msg.get("content", "")
                    else:
                        msg = getattr(choice, "message", None)
                        if msg:
                            try:
                                out = msg["content"]
                            except Exception:
                                out = getattr(msg, "content", None)
                except Exception:
                    out = None
                if out is None:
                    out = getattr(resp, "text", None) or ""
                return (out or "").strip()
            # fallback
            return ""
        except Exception as e:
            attempt += 1
            logger.exception("Error llamando a OpenAI (intento %s): %s", attempt, e)
            if attempt > max_retries:
                raise
            time.sleep(1 + attempt * 2)

def _openai_clean_text_to_markdown(text: str, instructions: Optional[str] = None, model: Optional[str] = None) -> str:
    """
    Envia texto crudo a OpenAI para convertir/limpiar en Markdown.
    """
    if not text or not text.strip():
        return ''
    model = model or OPENAI_MODEL_TEXT
    # build prompt
    base_instructions = (
        "Eres un asistente que convierte texto en Markdown limpio, legible y estructurado. "
        "Extrae títulos si existen, organiza en secciones y listas, corrige errores obvios, "
        "preserva tablas si las detectas (devuelve en Markdown), y responde únicamente con Markdown. "
        "No agregues explicaciones adicionales."
    )
    if instructions:
        base_instructions = instructions + " " + base_instructions

    messages = [
        {"role": "system", "content": base_instructions},
        {"role": "user", "content": f"Texto original:\n\n{text[:200000]}"}  # slice para no enviar exceso
    ]

    out = _call_openai_chat_completions(model=model, messages=messages)
    return text_to_md(out)

def txt_to_md(path: str) -> str:
    """Leer un archivo txt/md y devolver markdown limpio usando OpenAI para formateo."""
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            raw = f.read()
    except Exception as e:
        logger.exception("Error leyendo txt: %s", e)
        return ''
    # Si es corto, enviamos directamente; si es largo, chunk+resumir
    tokens = _estimate_tokens_from_chars(len(raw))
    if tokens > CHUNK_SIZE_TOKENS:
        chunks = _chunk_text_by_chars(raw, CHUNK_SIZE_TOKENS)
        summaries = []
        for i, c in enumerate(chunks):
            logger.info("txt_to_md: procesando chunk %s/%s", i+1, len(chunks))
            summaries.append(_openai_clean_text_to_markdown(c))
        # combinar y sintetizar
        combined = "\n\n".join(summaries)
        return _openai_clean_text_to_markdown(combined)
    else:
        return _openai_clean_text_to_markdown(raw)

def docx_to_md(path: str) -> str:
    """
    Extraer texto de docx y pasar a OpenAI para limpiar/convertir a Markdown.
    Si python-docx no está disponible, intentamos enviar bytes a OpenAI como fallback (menos recomendado).
    """
    text = ''
    if Document:
        try:
            doc = Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
            text = '\n\n'.join(paragraphs)
        except Exception as e:
            logger.exception("Error procesando docx localmente: %s", e)
            text = ''
    else:
        logger.warning("python-docx no instalado; no se puede extraer localmente docx.")

    if text:
        return _openai_clean_text_to_markdown(text)
    # Fallback: attempt to send bytes to OpenAI as a note (if supported). Here we try to extract raw bytes -> decode.
    try:
        with open(path, 'rb') as f:
            raw_bytes = f.read()
        # best-effort decode
        raw_text = raw_bytes.decode('utf-8', errors='ignore')
        if raw_text.strip():
            return _openai_clean_text_to_markdown(raw_text)
    except Exception:
        pass
    return ''

def pdf_to_md(path: str) -> str:
    """
    Procesamiento de PDF:
    1) Intentar extraer texto con pdfminer (rápido).
    2) Si no hay texto y pdf2image disponible -> renderizar páginas a imágenes y usar ocr_image_to_md (OpenAI Vision).
    3) Si todo falla y USE_LOCAL_OCR=True -> fallback a easyocr local.
    """
    text = ''
    if pdfminer_extract_text:
        try:
            text = pdfminer_extract_text(path)
            if text and text.strip():
                logger.info("pdf_to_md: texto extraido con pdfminer para %s", path)
                return _openai_clean_text_to_markdown(text)
        except Exception as e:
            logger.exception("pdf_to_md: fallo pdfminer: %s", e)
            text = ''

    # Si no hay texto, intentar renderizar a imagen por pagina y OCR vía OpenAI Vision
    if convert_from_path and Image:
        try:
            pages = convert_from_path(path)
            page_texts = []
            for i, page in enumerate(pages):
                try:
                    with io.BytesIO() as buf:
                        page.save(buf, format='PNG')
                        buf.seek(0)
                        page_bytes = buf.read()
                    page_md = ocr_image_to_md(page_bytes, lang=None)  # ocr_image_to_md maneja bytes
                    page_texts.append(page_md)
                except Exception as e:
                    logger.exception("pdf_to_md: fallo procesando pagina %s: %s", i, e)
            combined = "\n\n---\n\n".join([p for p in page_texts if p])
            if combined.strip():
                # sintetizar con OpenAI para obtener markdown coherente
                return _openai_clean_text_to_markdown(combined)
        except Exception as e:
            logger.exception("pdf_to_md: fallo render pdf->images: %s", e)

    # Fallback local OCR if configured
    if USE_LOCAL_OCR and easyocr:
        try:
            logger.info("pdf_to_md: intentando fallback local OCR con easyocr para %s", path)
            # Simple approach: attempt per page using pdf2image if available; else try easyocr on whole file path
            if convert_from_path and Image:
                pages = convert_from_path(path)
                page_texts = []
                for page in pages:
                    with io.BytesIO() as buf:
                        page.save(buf, format='PNG')
                        buf.seek(0)
                        page_bytes = buf.read()
                    # use easyocr directly on numpy array if available
                    try:
                        import numpy as np
                        img = Image.open(io.BytesIO(page_bytes)).convert("RGB")
                        arr = np.array(img)
                        reader = easyocr.Reader(['es', 'en'], gpu=False)
                        res = reader.readtext(arr, detail=0)
                        page_texts.append("\n".join(res))
                    except Exception:
                        logger.exception("pdf_to_md: easyocr per page failed")
                combined = "\n\n---\n\n".join(page_texts)
                return text_to_md(combined)
            else:
                # no pages rendering; last resort
                reader = easyocr.Reader(['es', 'en'], gpu=False)
                res = reader.readtext(path, detail=0)
                return text_to_md("\n".join(res))
        except Exception as e:
            logger.exception("pdf_to_md: fallo fallback easyocr: %s", e)

    # Si todo falla, devolver vacío
    logger.info("pdf_to_md: no se pudo extraer texto para %s", path)
    return ''

def _ocr_image_via_openai_bytes(image_bytes: bytes, ext: str = 'png', instructions: Optional[str] = None) -> str:
    """
    Envía una imagen en bytes a OpenAI Vision y solicita Markdown.
    """
    data_url = _image_file_to_data_url_bytes(image_bytes, ext)
    model = OPENAI_MODEL_VISION
    prompt = (
        (instructions or "")
        + " Extrae todo el texto legible de la imagen y devuélvelo en Markdown limpio. "
          "Corrige errores obvios de OCR y organiza el contenido en títulos, listas y párrafos cuando corresponda. "
          "Responde únicamente con Markdown."
    )

    messages = _build_image_messages(prompt, data_url)
    out = _call_openai_chat_completions(model=model, messages=messages)
    return text_to_md(out)

def ocr_image_to_md(file_obj, lang='es'):
    """
    Realiza OCR sobre una imagen y devuelve texto en formato Markdown.
    Compatible con objetos FieldFile, rutas, bytes, or file-likes.
    Implementación OpenAI-first: enviamos bytes a OpenAI Vision.
    """
    try:
        # Primero obtener bytes y extensión
        image_bytes = None
        ext = 'png'
        # 1: path string
        if isinstance(file_obj, str) and os.path.exists(file_obj):
            ext = os.path.splitext(file_obj)[1].lstrip('.')
            with open(file_obj, 'rb') as f:
                image_bytes = f.read()
        # 2: object with .path
        elif hasattr(file_obj, 'path') and os.path.exists(file_obj.path):
            ext = os.path.splitext(file_obj.path)[1].lstrip('.')
            with open(file_obj.path, 'rb') as f:
                image_bytes = f.read()
        # 3: file-like with open() method (Django FieldFile)
        elif hasattr(file_obj, 'open'):
            with file_obj.open('rb') as f:
                image_bytes = f.read()
        # 4: file-like with read()
        elif hasattr(file_obj, 'read'):
            # Might be an in-memory stream
            pos = None
            try:
                pos = file_obj.tell()
            except Exception:
                pos = None
            try:
                image_bytes = file_obj.read()
            finally:
                try:
                    if pos is not None:
                        file_obj.seek(pos)
                except Exception:
                    pass
        # 5: raw bytes
        elif isinstance(file_obj, (bytes, bytearray)):
            image_bytes = bytes(file_obj)
        else:
            raise ValueError(f"Tipo de archivo no soportado para OCR: {type(file_obj)}")

        if not image_bytes:
            return "[No se detectó contenido binario de imagen]"

        # Size guard
        size_mb = len(image_bytes) / (1024 * 1024)
        if size_mb > MAX_SEND_SIZE_MB:
            msg = f"[El archivo de imagen supera el límite de {MAX_SEND_SIZE_MB} MB ({size_mb:.2f} MB)]"
            logger.warning("ocr_image_to_md: %s", msg)
            return msg

        # Preferir OpenAI Vision
        try:
            return _ocr_image_via_openai_bytes(image_bytes, ext=ext, instructions=None)
        except Exception as e:
            logger.exception("ocr_image_to_md: fallo OpenAI Vision: %s", e)
            # fallback local?
            if USE_LOCAL_OCR and easyocr:
                try:
                    import numpy as np
                    img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                    arr = np.array(img)
                    reader = easyocr.Reader(['es', 'en'], gpu=False)
                    results = reader.readtext(arr, detail=0)
                    text = "\n".join(results).strip()
                    return text_to_md(text)
                except Exception as e2:
                    logger.exception("ocr_image_to_md: fallback easyocr falló: %s", e2)
                    return f"[Error OCR: {str(e)} | fallback failed: {str(e2)}]"
            else:
                return f"[Error OCR via OpenAI: {str(e)}]"

    except Exception as e:
        logger.exception("Error haciendo OCR a imagen con OpenAI/EasyOCR: %s", e)
        return f"[Error OCR: {str(e)}]"

def xlsx_to_md(path: str) -> str:
    """
    Extrae hojas de Excel a CSV por hoja y pide a OpenAI que sintetice insights y tabla en Markdown.
    Requiere openpyxl/pandas para extracción local. Si no está disponible, devuelve cadena vacía.
    """
    try:
        import pandas as pd
    except Exception:
        logger.warning("pandas/openpyxl no instalado; xlsx_to_md no disponible.")
        return ''

    try:
        xls = pd.read_excel(path, sheet_name=None)
        parts = []
        for sheet_name, df in xls.items():
            # small sample or full if small
            rows = df.head(10).to_csv(index=False)
            prompt = f"Hoja: {sheet_name}\nMuestra de filas:\n{rows}\n\nGenera un resumen breve y una tabla en Markdown con los 5 insights más importantes."
            md = _openai_clean_text_to_markdown(prompt)
            parts.append(f"## Hoja: {sheet_name}\n\n{md}")
        return "\n\n".join(parts)
    except Exception as e:
        logger.exception("xlsx_to_md error: %s", e)
        return ''

def pptx_to_md(path: str) -> str:
    """
    Extrae texto de pptx por diapositiva (si python-pptx está disponible) y pide resumen por slide + síntesis.
    """
    try:
        from pptx import Presentation
    except Exception:
        logger.warning("python-pptx no instalado; pptx_to_md no disponible.")
        return ''
    try:
        prs = Presentation(path)
        slides_md = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = shape.text.strip()
                    if t:
                        texts.append(t)
            slide_text = "\n\n".join(texts)
            if slide_text.strip():
                md_slide = _openai_clean_text_to_markdown(slide_text)
                slides_md.append(f"### Slide {i+1}\n\n{md_slide}")
        combined = "\n\n".join(slides_md)
        if combined.strip():
            return _openai_clean_text_to_markdown(combined)
        return ''
    except Exception as e:
        logger.exception("pptx_to_md error: %s", e)
        return ''